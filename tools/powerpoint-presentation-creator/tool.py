# /// script
# dependencies = [
#   "python-pptx",
#   "Pillow",
#   "requests",
# ]
# ///

from typing import Any, Optional, List, Dict, Callable
import json
import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
from shinkai_local_support import get_home_path
from shinkai_local_tools import shinkai_llm_prompt_processor

# CONFIG:
class CONFIG:
    pass

# INPUTS:
class INPUTS:
    instruction: str
    filename: Optional[str] = None

# OUTPUT:
class OUTPUT:
    saved_pptx_local_file_path: str
    status: str
    debug_info: str

# --- CONSTANTS ---
SLIDE_W = 10.0
SLIDE_H = 5.625

# Layout Spacing
MARGIN = 0.35
GAP = 0.15         # Used for both X and Y gaps in grids

# Vertical Anchors
# We give Title 1.25 height and start Content at 1.4 to be safe.
TITLE_Y = 0.2
TITLE_H = 1.25
CONTENT_Y = 1.4    

# Horizontal Anchors (40/60 Split)
TEXT_X = MARGIN
TEXT_W = 3.6
ZONE_X = 4.2
ZONE_W = SLIDE_W - ZONE_X - MARGIN
ZONE_Y = CONTENT_Y
ZONE_H = SLIDE_H - CONTENT_Y - MARGIN

STYLES = {
    "modern_clean": {
        "bg": (250, 250, 250), "title": (0, 0, 0), "text": (60, 60, 60), "accent": (0, 120, 215), "font": "Arial",
        "desc": "Standard corporate, clean, professional."
    },
    "dark_futuristic": {
        "bg": (15, 15, 20), "title": (0, 255, 200), "text": (220, 220, 220), "accent": (255, 0, 100), "font": "Consolas",
        "desc": "High-tech, AI, gaming, cyber-security."
    },
    "cinema_black": {
        "bg": (0, 0, 0), "title": (255, 255, 255), "text": (180, 180, 180), "accent": (229, 9, 20), "font": "Impact",
        "desc": "High contrast, media, movie scenarios, bold statements."
    },
    "midnight_blue": {
        "bg": (10, 25, 45), "title": (255, 255, 255), "text": (176, 196, 222), "accent": (255, 215, 0), "font": "Verdana",
        "desc": "Trustworthy, financial, academic."
    },
    "elegant_luxury": {
        "bg": (28, 28, 28), "title": (212, 175, 55), "text": (230, 230, 230), "accent": (169, 169, 169), "font": "Georgia",
        "desc": "Premium brands, history, literature."
    },
    "slate_developer": {
        "bg": (40, 44, 52), "title": (97, 175, 239), "text": (171, 178, 191), "accent": (198, 120, 221), "font": "Segoe UI",
        "desc": "Engineering, coding tutorials."
    },
    "forest_dark": {
        "bg": (20, 35, 25), "title": (152, 251, 152), "text": (240, 255, 240), "accent": (210, 180, 140), "font": "Trebuchet MS",
        "desc": "Nature, environment, sustainability."
    },
    "comic_book": {
        "bg": (15, 15, 15), "title": (255, 204, 0), "text": (255, 255, 255), "accent": (230, 40, 40), "font": "Comic Sans MS",
        "desc": "Comics, Storyboard, Playful, energetic, youth-oriented."
    },
}

def hex_to_rgb(rgb_tuple):
    return RGBColor(*rgb_tuple)

class ImageAnalysis:
    def __init__(self, path: str):
        self.path = path
        self.valid = False
        self.width = 0
        self.height = 0
        self.ratio = 1.0
        self.orientation = "square"
        try:
            if os.path.exists(path):
                with Image.open(path) as img:
                    self.valid = True
                    self.width, self.height = img.size
                    self.ratio = self.width / self.height if self.height != 0 else 1.0
                    if self.ratio > 1.25:
                        self.orientation = "wide"
                    elif self.ratio < 0.8:
                        self.orientation = "tall"
        except Exception:
            self.valid = False

def _sanitize_user_filename(name: Optional[str]) -> str:
    if not isinstance(name, str): return ""
    n = name.strip()
    if n == "": return ""
    n = os.path.splitext(n)[0]
    n = re.sub(r"[\s]+", " ", n)
    safe = "".join(c for c in n if c.isalnum() or c in (" ", "-", "_")).strip()
    return safe

# --- Layout Helpers ---
def calculate_text_geometry(images: List[str], home_path: str) -> str:
    valid_imgs = []
    for p in images:
        full_p = p if os.path.isabs(p) else os.path.join(home_path, p)
        a = ImageAnalysis(full_p)
        if a.valid:
            valid_imgs.append(a)
    c = len(valid_imgs)
    if c == 0:
        return "Space Available: FULL SLIDE (100% Width)."
    if c == 1:
        return "Space Available: FULL WIDTH (100%), but only TOP 30% Height." if valid_imgs[0].orientation == "wide" else "Space Available: LEFT 50% Width, FULL Height."
    return "Space Available: LEFT 40% Width, FULL Height. (Standard bullets)."

def smart_place_image(slide, img_obj: ImageAnalysis, x, y, w, h, border_color=None):
    if not img_obj.valid: return
    box_ratio = w / h
    if img_obj.ratio > box_ratio:
        render_w = w
        render_h = w / img_obj.ratio
    else:
        render_h = h
        render_w = h * img_obj.ratio
    off_x = x + (w - render_w) / 2
    off_y = y + (h - render_h) / 2
    try:
        pic = slide.shapes.add_picture(img_obj.path, Inches(off_x), Inches(off_y), width=Inches(render_w), height=Inches(render_h))
        if border_color:
            pic.line.color.rgb = border_color
            pic.line.width = Pt(1.5)
    except Exception:
        pass

# --- LayoutEngine ---
class LayoutEngine:
    def __init__(self, prs: Presentation, style: Dict[str, Any], home_path: str):
        self.prs = prs
        self.style = style
        self.home_path = home_path

    def _setup_base_slide(self, title: str):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(SLIDE_W), Inches(SLIDE_H))
        bg.fill.solid()
        bg.fill.fore_color.rgb = hex_to_rgb(self.style["bg"])
        bg.line.fill.background()
        
        # Title Box - Taller for 2 lines
        t_box = slide.shapes.add_textbox(Inches(MARGIN), Inches(TITLE_Y), Inches(SLIDE_W - MARGIN*2), Inches(TITLE_H))
        t_box.text_frame.text = str(title) if title else ""
        if t_box.text_frame.paragraphs:
            p = t_box.text_frame.paragraphs[0]
            p.font.name = self.style["font"]
            p.font.size = Pt(30)
            p.font.bold = True
            p.font.color.rgb = hex_to_rgb(self.style["title"])
        return slide

    def _add_text_block(self, slide, content, x, y, w, h):
        if not content:
            return
        content_list = content if isinstance(content, list) else [str(content)]
        content_list = [str(c) for c in content_list if c]
        if not content_list: return
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.word_wrap = True
        for i, line in enumerate(content_list):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.font.name = self.style["font"]
            p.font.size = Pt(18)
            p.font.color.rgb = hex_to_rgb(self.style["text"])
            p.space_after = Pt(10)

    def layout_single(self, slide_data, img: ImageAnalysis):
        slide = self._setup_base_slide(slide_data.get("title", ""))
        has_text = bool(slide_data.get("content"))
        if has_text and img.orientation == "wide":
            self._add_text_block(slide, slide_data["content"], MARGIN, CONTENT_Y, SLIDE_W - MARGIN*2, 1.2)
            img_y = CONTENT_Y + 1.5
            img_h = SLIDE_H - img_y - MARGIN
            smart_place_image(slide, img, MARGIN, img_y, SLIDE_W - MARGIN*2, img_h, hex_to_rgb(self.style["accent"]))
        elif has_text:
            self._add_text_block(slide, slide_data["content"], MARGIN, CONTENT_Y, 4.5, 4.0)
            smart_place_image(slide, img, 5.0, CONTENT_Y, 4.6, 3.8, hex_to_rgb(self.style["accent"]))
        else:
            smart_place_image(slide, img, MARGIN, CONTENT_Y, SLIDE_W - MARGIN*2, SLIDE_H - CONTENT_Y - MARGIN, hex_to_rgb(self.style["accent"]))

    def layout_dual(self, slide_data, imgs: List[ImageAnalysis]):
        slide = self._setup_base_slide(slide_data.get("title", ""))
        has_text = bool(slide_data.get("content"))
        
        # Geometry Logic
        use_text_x = TEXT_X
        use_text_w = TEXT_W
        use_zone_x = ZONE_X
        use_zone_w = ZONE_W
        
        # Check for "Stacked" scenario (Wide images stacked on right)
        # Standard Zone H is ~3.8. Two images means each is ~1.8h.
        # 16:9 images at 1.8h are only ~3.2w.
        # Current Zone W is 5.45. That leaves >2.0 inches of whitespace.
        # We can safely give that space to the text!
        avg_ratio = sum([i.ratio for i in imgs]) / 2
        is_stacked = avg_ratio > 1.2
        
        if is_stacked and has_text:
            # Shift geometry for Stacked Dual Images
            # New Text Width: 5.5 (Leaves 4.15 for images)
            use_text_w = 5.5
            use_zone_x = MARGIN + use_text_w + 0.25 # 0.25 gutter
            use_zone_w = SLIDE_W - use_zone_x - MARGIN
        
        if has_text:
            self._add_text_block(slide, slide_data["content"], use_text_x, CONTENT_Y, use_text_w, 4.0)
        else:
            use_zone_x, use_zone_w = MARGIN, SLIDE_W - MARGIN*2

        zy, zh = ZONE_Y, ZONE_H

        if is_stacked:
            # Stacked Layout
            h_half = (zh - GAP) / 2
            smart_place_image(slide, imgs[0], use_zone_x, zy, use_zone_w, h_half, hex_to_rgb(self.style["accent"]))
            smart_place_image(slide, imgs[1], use_zone_x, zy + h_half + GAP, use_zone_w, h_half, hex_to_rgb(self.style["accent"]))
        else:
            # Side-by-Side Layout
            w_half = (use_zone_w - GAP) / 2
            smart_place_image(slide, imgs[0], use_zone_x, zy, w_half, zh, hex_to_rgb(self.style["accent"]))
            smart_place_image(slide, imgs[1], use_zone_x + w_half + GAP, zy, w_half, zh, hex_to_rgb(self.style["accent"]))

    def layout_smart_3(self, slide_data, imgs: List[ImageAnalysis]):
        slide = self._setup_base_slide(slide_data.get("title", ""))
        has_text = bool(slide_data.get("content"))
        if has_text:
            self._add_text_block(slide, slide_data["content"], TEXT_X, CONTENT_Y, TEXT_W, 4.0)
            zx, zw = ZONE_X, ZONE_W
        else:
            zx, zw = MARGIN, SLIDE_W - MARGIN*2
        zy, zh = ZONE_Y, ZONE_H
        wides = [img for img in imgs if img.orientation == "wide"]
        talls = [img for img in imgs if img.orientation == "tall"]
        if len(wides) >= 1:
            anchor = wides[0]
            remaining = [img for img in imgs if img != anchor]
            h_half = (zh - GAP) / 2
            w_half = (zw - GAP) / 2
            smart_place_image(slide, remaining[0], zx, zy, w_half, h_half, hex_to_rgb(self.style["accent"]))
            smart_place_image(slide, remaining[1], zx + w_half + GAP, zy, w_half, h_half, hex_to_rgb(self.style["accent"]))
            smart_place_image(slide, anchor, zx, zy + h_half + GAP, zw, h_half, hex_to_rgb(self.style["accent"]))
        elif len(talls) >= 1:
            anchor = talls[0]
            remaining = [img for img in imgs if img != anchor]
            w_big = (zw - GAP) / 2
            w_small = w_big
            h_small = (zh - GAP) / 2
            smart_place_image(slide, anchor, zx, zy, w_big, zh, hex_to_rgb(self.style["accent"]))
            smart_place_image(slide, remaining[0], zx + w_big + GAP, zy, w_small, h_small, hex_to_rgb(self.style["accent"]))
            smart_place_image(slide, remaining[1], zx + w_big + GAP, zy + h_small + GAP, w_small, h_small, hex_to_rgb(self.style["accent"]))
        else:
            h_half = (zh - GAP) / 2
            w_half = (zw - GAP) / 2
            smart_place_image(slide, imgs[0], zx, zy, w_half, h_half, hex_to_rgb(self.style["accent"]))
            smart_place_image(slide, imgs[1], zx + w_half + GAP, zy, w_half, h_half, hex_to_rgb(self.style["accent"]))
            smart_place_image(slide, imgs[2], zx, zy + h_half + GAP, zw, h_half, hex_to_rgb(self.style["accent"]))

    def layout_smart_4(self, slide_data, imgs: List[ImageAnalysis]):
        slide = self._setup_base_slide(slide_data.get("title", ""))
        has_text = bool(slide_data.get("content"))
        if has_text:
            self._add_text_block(slide, slide_data["content"], TEXT_X, CONTENT_Y, TEXT_W, 4.0)
            zx, zw = ZONE_X, ZONE_W
        else:
            zx, zw = MARGIN, SLIDE_W - MARGIN*2
        zy, zh = ZONE_Y, ZONE_H
        wides = [img for img in imgs if img.orientation == "wide"]
        sorted_imgs = []
        if len(wides) >= 2:
            others = [img for img in imgs if img not in wides]
            sorted_imgs = others + wides
            if len(sorted_imgs) > 4: sorted_imgs = sorted_imgs[:4]
        else:
            sorted_imgs = imgs[:4]
        
        # GAP ensures vertical and horizontal spacing are identical
        w_half = (zw - GAP) / 2
        h_half = (zh - GAP) / 2
        
        coords = [
            (zx, zy), 
            (zx + w_half + GAP, zy), 
            (zx, zy + h_half + GAP), 
            (zx + w_half + GAP, zy + h_half + GAP)
        ]
        
        for i in range(min(4, len(sorted_imgs))):
            smart_place_image(slide, sorted_imgs[i], coords[i][0], coords[i][1], w_half, h_half, hex_to_rgb(self.style["accent"]))

    def render_slide(self, slide_data):
        raw_imgs = slide_data.get("images", [])
        if not isinstance(raw_imgs, list): raw_imgs = []
        valid_imgs = []
        for p in raw_imgs:
            if not p: continue
            fp = p if os.path.isabs(p) else os.path.join(self.home_path, p)
            a = ImageAnalysis(fp)
            if a.valid: valid_imgs.append(a)
        
        cnt = len(valid_imgs)
        if cnt == 0:
            slide = self._setup_base_slide(slide_data.get("title", ""))
            self._add_text_block(slide, slide_data.get("content", []), MARGIN, CONTENT_Y, SLIDE_W - MARGIN*2, 4.0)
        elif cnt == 1:
            self.layout_single(slide_data, valid_imgs[0])
        elif cnt == 2:
            self.layout_dual(slide_data, valid_imgs)
        elif cnt == 3:
            self.layout_smart_3(slide_data, valid_imgs)
        else:
            self.layout_smart_4(slide_data, valid_imgs)

# --- LLM interaction with retry logic ---
async def _call_llm_with_retry(prompt: str, validator_fn: Optional[Callable[[Dict], None]] = None, max_retries: int = 1) -> Dict[str, Any]:
    """
    Calls LLM, parses JSON, and optionally validates content using validator_fn.
    validator_fn should raise ValueError if the content is invalid.
    """
    attempt = 0
    last_raw = None
    while attempt <= max_retries:
        attempt += 1
        try:
            res = await shinkai_llm_prompt_processor({"prompt": prompt, "format": "json"})
            raw = res.get("message", "{}")
            last_raw = raw
            cleaned = raw.strip().replace("```json", "").replace("```", "")
            parsed = json.loads(cleaned)
            
            # 1. Structural Validation
            if not isinstance(parsed, dict):
                raise ValueError("Output must be a JSON object (dictionary).")
            
            # 2. Logic/Schema Validation
            if validator_fn:
                validator_fn(parsed)

            return {"ok": True, "data": parsed, "raw": raw, "attempts": attempt}
        
        except Exception as e:
            # Prepare guidance for retry
            guidance = (
                f"Previous output was invalid. Issue: {e}. "
                f"Returned content (truncated): {str(last_raw)[:500]!s}. "
                "Please fix the error and return valid JSON."
            )
            prompt = prompt + "\n\n" + guidance
            if attempt > max_retries:
                return {"ok": False, "data": last_raw if last_raw is not None else "", "raw": last_raw, "error": str(e), "attempts": attempt}

async def run(config: CONFIG, inputs: INPUTS) -> OUTPUT:
    output = OUTPUT()
    output.saved_pptx_local_file_path = ""
    output.status = "error"
    output.debug_info = ""

    # Validate inputs
    if not inputs or not hasattr(inputs, "instruction") or not isinstance(inputs.instruction, str) or inputs.instruction.strip() == "":
        output.debug_info = "Input 'instruction' is required and must be a non-empty string."
        return output

    try:
        home_path = await get_home_path()
    except Exception as e:
        output.debug_info = f"Failed to determine home path: {e}"
        return output

    style_list = ", ".join(STYLES.keys())
    style_options = {k: v["desc"] for k, v in STYLES.items()}

    # --- Schema Validation Logic ---
    def validate_schema_strict(data: Dict[str, Any], require_style: bool):
        # Top level check
        if require_style and "detected_style" not in data:
            raise ValueError("Missing required key: 'detected_style'")
        if "slides" not in data:
            raise ValueError("Missing required key: 'slides'")
        
        slides = data["slides"]
        if not isinstance(slides, list):
            raise ValueError("'slides' must be a list.")
        
        ALLOWED_KEYS = {"title", "content", "images"}
        
        for i, s in enumerate(slides):
            if not isinstance(s, dict):
                raise ValueError(f"Slide {i} must be a JSON object.")
            
            # strict key check
            unknown = set(s.keys()) - ALLOWED_KEYS
            if unknown:
                raise ValueError(f"Slide {i} contains forbidden keys: {list(unknown)}. Only allowed: {list(ALLOWED_KEYS)}.")
            
            # image path check
            if "images" in s:
                imgs = s["images"]
                if not isinstance(imgs, list):
                    raise ValueError(f"Slide {i} 'images' must be a list.")
                for path in imgs:
                    full_p = path if os.path.isabs(path) else os.path.join(home_path, path)
                    if not os.path.exists(full_p):
                        raise ValueError(f"Slide {i} references missing local file: '{path}'. Check path or remove image.")

    # Pass 1: Architect
    prompt_1 = f"""
You are a Slides Presentation Architect.

Goal: Convert user request into a JSON object for slide generation.

Step 1: Style Selection
Review the `Available Styles` JSON below and select the key that best matches the user's intent.
Available Styles:
{json.dumps(style_options, indent=2)}

Step 2: Content Generation
- Extract a title for each slide (very short, must fit in 2 lines max with ultre large font size, from 1 to 11 words max).
- Extract clear, qualitative content. 
- Identify ALL local image paths mentioned.
You are writing for a visual medium (16:9 Slide).
Do NOT "summarize" generically. Instead, edit the text for information quality, impact and layout fit. Keep qualitative meaning but remove fluff.
Keep in mind that large blocks of text are annoying to read on screen, and that long vertical blocks of text would be off the slide.

Rules for Text Density:
- **0 Images:** Full slide available. You might write more.
- **1 or more Images:** You will have less space to write, so content should be shorter.

JSON Schema:
{{
  "detected_style": "STYLE_NAME",
  "slides": [
    {{ "title": "My Title", "content": ["Point 1", "Point 2"], "images": ["/path/1.jpg"] }}
  ]
}}
User Request: {inputs.instruction}
"""
    # Validator lambda for Pass 1
    p1_validator = lambda d: validate_schema_strict(d, require_style=True)
    
    p1_result = await _call_llm_with_retry(prompt_1, validator_fn=p1_validator, max_retries=1)
    if not p1_result.get("ok"):
        output.debug_info = f"Pass 1 failed after {p1_result.get('attempts')} attempts. Error: {p1_result.get('error')}"
        return output

    draft_data = p1_result["data"]
    slides_draft = draft_data.get("slides", [])

    # Enrich geometry
    for slide in slides_draft:
        imgs = slide.get("images", [])
        if not isinstance(imgs, list): imgs = []
        slide["layout_geometry_hint"] = calculate_text_geometry(imgs, home_path)

    # Pass 2: Typesetter
    prompt_2 = f"""
You are a Slides Editor (Typesetter).
Goal: Adjust the slide content to fit the specific geometric space available on the screen.

Input Data (Draft Content + Geometry Hint):
{json.dumps(slides_draft, indent=2)}

Instructions:
1. Iterate through every slide.
2. Read the "layout_geometry_hint". This describes your available text box (e.g., "Left 50% Width" or "Top 30% Height").
3. REWRITE the "content" to fit comfortably in that space.
   - If space is "Full Width, Top 30% Height": You have a wide but short strip. Use only a few long lines.
   - If space is "Left 40% Width": You have a narrow column. Use shorter lines/bullets.
   - If space is "Full Slide": You have plenty of room.
4. Keep the qualitative meaning. Just format/resize the text logic.

Return JSON Schema:
{{
  "slides": [
    {{ "title": "My Title", "content": ["Point 1", "Point 2"], "images": ["/path/1.jpg"] }}
  ]
}}

Important: Return ONLY valid JSON.
"""
    # Validator lambda for Pass 2 (style not required here)
    p2_validator = lambda d: validate_schema_strict(d, require_style=False)
    
    p2_result = await _call_llm_with_retry(prompt_2, validator_fn=p2_validator, max_retries=1)
    
    if not p2_result.get("ok"):
        # Fallback to Pass 1
        err_msg = p2_result.get("error", "Validation error")
        output.debug_info = f"Pass 2 failed ({err_msg}). Falling back to Pass 1."
        final_slides = slides_draft
    else:
        final_slides = p2_result["data"].get("slides", slides_draft)

    # Render
    try:
        prs = Presentation()
        prs.slide_width = Inches(SLIDE_W)
        prs.slide_height = Inches(SLIDE_H)
    except Exception as e:
        output.debug_info = f"PPTX Init failed: {e}"
        return output

    detected_style = draft_data.get("detected_style", "modern_clean")
    style_config = STYLES.get(detected_style, STYLES["modern_clean"])
    engine = LayoutEngine(prs, style_config, home_path)

    render_errors = []
    for idx, sd in enumerate(final_slides):
        try:
            # Restore images if Pass 2 dropped them, though validator checks this too
            if "images" not in sd and idx < len(slides_draft):
                sd["images"] = slides_draft[idx].get("images", [])
            engine.render_slide(sd)
        except Exception as e:
            render_errors.append(f"Slide {idx}: {e}")

    # Save
    user_fname = getattr(inputs, "filename", None)
    safe_base = _sanitize_user_filename(user_fname) if user_fname else ""
    if not safe_base:
        safe_base = "presentation"
    
    final_filename = safe_base + ".pptx"
    final_path = os.path.join(home_path, final_filename)
    
    ctr = 1
    root, ext = os.path.splitext(final_path)
    while os.path.exists(final_path):
        final_path = f"{root}_{ctr}{ext}"
        ctr += 1

    try:
        prs.save(final_path)
        output.saved_pptx_local_file_path = final_path
        output.status = "success"
        info = f"Generated {len(final_slides)} slides. Style: '{detected_style}'."
        if output.debug_info: info += f" {output.debug_info}"
        if render_errors: info += f" Errors: {render_errors}"
        output.debug_info = info
    except Exception as e:
        output.saved_pptx_local_file_path = ""
        output.status = "error"
        output.debug_info = f"Save failed: {e}"

    return output