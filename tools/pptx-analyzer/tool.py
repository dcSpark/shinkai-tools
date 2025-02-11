# /// script
# dependencies = [
#   "requests",
#   "python-pptx",
# ]
# ///

from typing import Any, Optional, List, Dict
import os
import requests
from pptx import Presentation
from shinkai_local_tools import shinkai_llm_prompt_processor
from shinkai_local_support import get_home_path 

class CONFIG:
    pass

class INPUTS:
    operation: str
    file_path: Optional[str] = None
    prompt: Optional[str] = None

class OUTPUT:
    content: Optional[List[str]] = None
    analysis: Optional[str] = None
    status: Optional[str] = None

async def _read_presentation(file_path: str, working_file_path: str) -> OUTPUT:
    try:
        if file_path.startswith('http://') or file_path.startswith('https://'):
            response = requests.get(file_path)
            with open(working_file_path, 'wb') as f:
                f.write(response.content)
        else:
            with open(file_path, 'rb') as f:
                with open(working_file_path, 'wb') as w:
                    w.write(f.read())
        o = OUTPUT()
        o.status = "success"
        return o
    except Exception as e:
        o = OUTPUT()
        o.status = "error"
        o.analysis = f"Failed to read presentation: {str(e)}"
        return o

async def _retrieve_presentation(working_file_path: str) -> OUTPUT:
    if not os.path.exists(working_file_path):
        return _create_error_output("No presentation file loaded. Please use READ operation first.")
        
    try:
        # Read the file content
        with open(working_file_path, 'rb') as f:
            content = f.read()
        
        # Remove the file
        os.remove(working_file_path)
        
        # Rewrite the file
        with open(working_file_path, 'wb') as f:
            f.write(content)
            
        prs = Presentation(working_file_path)
        text_content = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_content.append(run.text)
        o = OUTPUT()
        o.content = text_content
        o.status = "success"
        return o
    except Exception as e:
        return _create_error_output(f"Failed to read presentation: {str(e)}")

async def _analyze_presentation(working_file_path: str, prompt: Optional[str]) -> OUTPUT:
    if not os.path.exists(working_file_path):
        return _create_error_output("No presentation file loaded. Please use READ operation first.")
        
    prs = Presentation(working_file_path)
    content = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    content.append(run.text)
    analysis_prompt = prompt if prompt else "Summarize the contents"
    analysis_result = await shinkai_llm_prompt_processor({
        "prompt": f"The user has asked a question about the following presentation <content>{' '.join(content)}</content> <question>{analysis_prompt}</question>", 
        "format": "text"
    })
    o = OUTPUT()
    o.content = content
    o.analysis = analysis_result['message']
    o.status = "success"
    return o

async def _update_presentation(working_file_path: str, update_request: Dict[str, Any]) -> OUTPUT:
    if not os.path.exists(working_file_path):
        return _create_error_output("No presentation file loaded. Please use READ operation first.")
        
    try:
        prs = Presentation(working_file_path)
        # Here you would implement the logic to update the presentation based on update_request
        # Example: prs.slides[0].shapes.title.text = update_request.get("title", prs.slides[0].shapes.title.text)
        prs.save(working_file_path)  # Save the updated presentation
        o = OUTPUT()
        o.status = "success"
        return o
    except Exception as e:
        o = OUTPUT()
        o.status = "error"
        o.analysis = f"Failed to update presentation: {str(e)}"
        return o

def _create_error_output(message: str) -> OUTPUT:
    o = OUTPUT()
    o.analysis = message
    o.status = "error"
    return o

async def run(config: CONFIG, inputs: INPUTS) -> OUTPUT:
    home_path = await get_home_path()
    working_file_path = os.path.join(home_path, 'working.pptx')
    
    # Convert operation to uppercase for case-insensitive comparison
    operation = inputs.operation.upper()
    
    # Always read presentation if file_path is provided
    if inputs.file_path:
        read_result = await _read_presentation(inputs.file_path, working_file_path)
        if read_result.status == "error":
            return read_result
    
    if operation == "READ":
        return read_result if inputs.file_path else _create_error_output("No file path provided for READ operation")
    elif operation == "RETRIEVE":
        return await _retrieve_presentation(working_file_path)
    elif operation == "ANALYZE":
        return await _analyze_presentation(working_file_path, inputs.prompt)
    elif operation == "UPDATE":
        return await _update_presentation(working_file_path, inputs.prompt)

    return _create_error_output(f"Invalid operation: {inputs.operation}")